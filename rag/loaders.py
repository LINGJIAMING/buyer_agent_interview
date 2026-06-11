# -*- coding: utf-8 -*-
"""PDF / PPTX / DOCX 文本抽取。"""
from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import List


@dataclass
class RawBlock:
    text: str
    page_or_slide: str
    section_title: str = ""


def load_pdf(path: Path) -> List[RawBlock]:
    import fitz  # PyMuPDF

    blocks: List[RawBlock] = []
    doc = fitz.open(path)
    for i, page in enumerate(doc, start=1):
        text = page.get_text("text") or ""
        text = text.strip()
        if len(text) < 20:
            continue
        blocks.append(RawBlock(text=text, page_or_slide=str(i)))
    doc.close()
    return blocks


def load_pptx(path: Path) -> List[RawBlock]:
    from pptx import Presentation

    blocks: List[RawBlock] = []
    prs = Presentation(str(path))
    for i, slide in enumerate(prs.slides, start=1):
        parts: List[str] = []
        title = ""
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            t = (shape.text or "").strip()
            if not t:
                continue
            if shape == slide.shapes.title or (not title and len(t) < 80):
                title = t
            parts.append(t)
        text = "\n".join(parts).strip()
        if len(text) < 15:
            continue
        blocks.append(RawBlock(text=text, page_or_slide=str(i), section_title=title))
    return blocks


def load_docx(path: Path) -> List[RawBlock]:
    from docx import Document

    blocks: List[RawBlock] = []
    doc = Document(str(path))
    buf: List[str] = []
    section = ""

    def flush():
        nonlocal buf, section
        text = "\n".join(buf).strip()
        if len(text) >= 20:
            blocks.append(RawBlock(text=text, page_or_slide="doc", section_title=section))
        buf = []

    for para in doc.paragraphs:
        t = (para.text or "").strip()
        if not t:
            flush()
            continue
        style = (para.style.name if para.style else "") or ""
        if "Heading" in style or t.endswith("：") and len(t) < 40:
            flush()
            section = t
        buf.append(t)
    flush()
    return blocks


def load_document(path: Path) -> List[RawBlock]:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return load_pdf(path)
    if ext == ".pptx":
        return load_pptx(path)
    if ext == ".docx":
        return load_docx(path)
    raise ValueError(f"unsupported: {path}")
